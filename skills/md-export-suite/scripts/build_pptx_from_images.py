#!/usr/bin/env python3
"""Assemble sorted slide images into a PPTX deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ASPECTS = {
    "16:9": (13.333333, 7.5),
    "4:5": (10.0, 12.5),
    "1:1": (10.0, 10.0),
    "9:16": (7.5, 13.333333),
}


def parse_size(aspect: str) -> tuple[float, float]:
    if aspect in ASPECTS:
        return ASPECTS[aspect]
    if "x" in aspect.lower():
        w, h = aspect.lower().split("x", 1)
        return float(w), float(h)
    raise ValueError(f"Unsupported aspect: {aspect}")


def build(image_dir: Path, output: Path, aspect: str, pattern: str) -> None:
    images = sorted(image_dir.glob(pattern))
    if not images:
        raise SystemExit(f"No images found in {image_dir} matching {pattern}")

    width, height = parse_size(aspect)
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)
    blank = prs.slide_layouts[6]

    for image in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=prs.slide_width, height=prs.slide_height)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(output)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build PPTX from slide images.")
    parser.add_argument("image_dir", help="Directory with slide images.")
    parser.add_argument("--output", "-o", required=True, help="Output PPTX path.")
    parser.add_argument("--aspect", default="16:9", help="16:9, 4:5, 1:1, 9:16, or WxH inches.")
    parser.add_argument("--pattern", default="slide-*.png", help="Image glob pattern.")
    args = parser.parse_args()
    build(Path(args.image_dir), Path(args.output), args.aspect, args.pattern)


if __name__ == "__main__":
    main()
